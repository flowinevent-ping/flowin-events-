# -*- coding: utf-8 -*-
import subprocess, os, io, sys, numpy as np
sys.path.insert(0,"/home/claude/vid")
import nds_lib as L
from pptx import Presentation
from pptx.oxml.ns import qn
from PIL import Image
from pyzbar.pyzbar import decode
KD="/home/claude/flowin/admin/public/nds/kit-digital"
QDIR="/home/claude/flowin/admin/public/nds/visuels-src/qr"
WORK="/home/claude/pptx"
FLOWIN=Image.open("/home/claude/pptx/flowin_logo.png").convert("RGBA")
PARTNERS=["bergerie","pegase","utile","nook","giordano","charvolin","carrosserie-gp"]

def png_bytes(im):
    o=io.BytesIO(); im.save(o,"PNG"); return o.getvalue()
def get_blob_shape(sl,name): return [s for s in sl.shapes if s.name==name][0]
def set_blob(shape,b):
    rId=shape._element.xpath(".//a:blip")[0].get(qn("r:embed"))
    shape.part.related_part(rId)._blob=b
def clean_ticket_blob(shape):
    im=Image.open(io.BytesIO(shape.image.blob)).convert("RGBA"); a=np.array(im)
    a[a[:,:,3]<130,3]=0
    return png_bytes(Image.fromarray(a))

def build_pptx(slug):
    if os.path.exists(f"{KD}/{slug}/nds_a4_{slug}-editable.pptx"):
        p=Presentation(f"{KD}/{slug}/nds_a4_{slug}-editable.pptx")
    else:  # nook : reconstruire depuis bergerie
        p=Presentation(f"{KD}/bergerie/nds_a4_bergerie-editable.pptx")
        sl=p.slides[0]
        set_blob(get_blob_shape(sl,"Image 3"), png_bytes(L.logo_card(f"/home/claude/repo/admin/public/nds/partenaires/{slug}.png",420,420)))
        qsrc=f"{QDIR}/ev-nds-{slug}.png"
        if not os.path.exists(qsrc): qsrc=f"{KD}/{slug}/qr-station-{slug}.png"
        qim=Image.open(qsrc).convert("RGB").resize((1000,1000),Image.NEAREST)
        set_blob(get_blob_shape(sl,"Image 6"), png_bytes(qim))
    sl=p.slides[0]
    # 1) encadré transparent des billets retiré
    for nm in ("Image 4","Image 5"):
        sh=get_blob_shape(sl,nm); set_blob(sh, clean_ticket_blob(sh))
    # 2) STATION JEUX -> STATION JEU
    t3=get_blob_shape(sl,"Text 3")
    for para in t3.text_frame.paragraphs:
        for r in para.runs:
            if "JEUX" in r.text.upper(): r.text=r.text.upper().replace("JEUX","JEU")
    out=f"{WORK}/pa_{slug}.pptx"; p.save(out); return out

def footer_fix(png):
    a=np.array(png.convert("RGB")); H,W,_=a.shape
    y0=int(27.6/29.7*H)
    for y in range(y0,H): a[y,:,:]=a[y,8,:]   # retire mail + trait
    im=Image.fromarray(a).convert("RGBA")
    fw=int(W*0.26); fh=int(FLOWIN.height*fw/FLOWIN.width)
    im.alpha_composite(FLOWIN.resize((fw,fh),Image.LANCZOS),(int(W/2-fw/2),int(28.55/29.7*H-fh/2)))
    return im.convert("RGB")

for s in PARTNERS:
    pptx=build_pptx(s)
    subprocess.run(["python3","/mnt/skills/public/pptx/scripts/office/soffice.py","--headless","--convert-to","pdf","--outdir",WORK,pptx],check=False,capture_output=True)
    subprocess.run(["pdftoppm","-png","-r","300","-singlefile",f"{WORK}/pa_{s}.pdf",f"{WORK}/pa_{s}"],check=False)
    im=footer_fix(Image.open(f"{WORK}/pa_{s}.png"))
    im.save(f"{KD}/{s}/nds_a4_{s}.pdf","PDF",resolution=300.0)
    im.resize((int(29.7/2.54*200),int(42/2.54*200)),Image.LANCZOS).save(f"{KD}/{s}/nds_a3_{s}.pdf","PDF",resolution=200.0)
    im.resize((int(32/2.54*200),int(45/2.54*200)),Image.LANCZOS).save(f"{KD}/{s}/affiche-partenaire_32x45_{s}.pdf","PDF",resolution=200.0)
    im.resize((im.width//2,im.height//2)).save(f"{KD}/{s}/nds_a3_{s}.png")
    r=decode(im); qr=r[0].data.decode().split("parcours/")[-1] if r else "NO"
    import pytesseract
    foot='flowinevent' in pytesseract.image_to_string(im.crop((0,int(27.6/29.7*im.height),im.width,im.height))).lower()
    print(f"{s:16} QR={qr}  mail={foot}")
print("DONE")
