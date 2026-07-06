# -*- coding: utf-8 -*-
import os, subprocess, io, copy
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
from pyzbar.pyzbar import decode
REF="/home/claude/pptx/ref.pptx"; QDIR="/home/claude/flowin/admin/public/nds/visuels-src/qr"
ND="/home/claude/flowin/admin/public/nds/kit-digital/nds"; WORK="/home/claude/pptx"
STAMAP={"caisse-1":"CAISSE 1","caisse-2":"CAISSE 2","bar-1":"BAR 1","bar-2":"BAR 2","bar-3":"BAR 3"}

LOGO_NAMES={"Picture 12","Picture 13","Picture 14","Picture 15","Picture 16"}
def add_card_behind(sl, pic, padx, pady):
    card=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pic.left-padx, pic.top-pady, pic.width+2*padx, pic.height+2*pady)
    card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF)
    card.line.fill.background()
    try: card.shadow.inherit=False
    except Exception: pass
    # arrondi doux
    try: card.adjustments[0]=0.18
    except Exception: pass
    el=card._element; el.getparent().remove(el); pic._element.addprevious(el)

def adjust(p):
    sl=p.slides[0]
    # STATION JEUX -> STATION JEU (singulier, affiches)
    for s in sl.shapes:
        if s.name=="Rounded Rectangle 4":
            for para in s.text_frame.paragraphs:
                for r in para.runs:
                    if "JEUX" in r.text.upper(): r.text=r.text.upper().replace("JEUX","JEU")
    # descendre NOS PARTENAIRES (mord sur JOUE)
    for s in sl.shapes:
        if s.name=="TextBox 10":
            s.top=Cm(21.9)
    return p

# --- template A4 ajuste ---
pA4=adjust(Presentation(REF)); pA4.save(f"{WORK}/tplA4.pptx")
# --- template 32x45 ajuste (scale) ---
NEW_W=int(round(32*360000)); NEW_H=int(round(45*360000))
p=Presentation(f"{WORK}/tplA4.pptx"); OW,OH=p.slide_width,p.slide_height
sx=NEW_W/OW; sy=NEW_H/OH; p.slide_width=NEW_W; p.slide_height=NEW_H
for s in p.slides:
    for sh in s.shapes:
        sh.left=int(sh.left*sx); sh.top=int(sh.top*sy); sh.width=int(sh.width*sx); sh.height=int(sh.height*sy)
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size is not None: r.font.size=Pt(round(r.font.size.pt*sx,1))
p.save(f"{WORK}/tpl3245.pptx")

def hi_qr(st):
    im=Image.open(f"{QDIR}/ev-nds-{st}.png").convert("RGB").resize((1200,1200),Image.NEAREST)
    b=io.BytesIO(); im.save(b,"PNG"); return b.getvalue()

def build(tpl,fmt):
    for st,label in STAMAP.items():
        p=Presentation(tpl); sl=p.slides[0]
        q=[s for s in sl.shapes if s.name=="Picture 8"][0]
        rId=q._element.xpath(".//a:blip")[0].get(qn("r:embed")); q.part.related_part(rId)._blob=hi_qr(st)
        c=[s for s in sl.shapes if s.name=="Rounded Rectangle 3"][0]; runs=c.text_frame.paragraphs[0].runs
        if runs:
            runs[0].text=label
            for r in runs[1:]: r.text=""
        outp=f"{ND}/station-jeu_{fmt}_{st}_editable.pptx"; p.save(outp)
        subprocess.run(["python3","/mnt/skills/public/pptx/scripts/office/soffice.py","--headless","--convert-to","pdf","--outdir",WORK,outp],check=False,capture_output=True)
        src=f"{WORK}/station-jeu_{fmt}_{st}_editable.pdf"; dst=f"{ND}/station-jeu_{fmt}_{st}.pdf"
        if os.path.exists(src): os.replace(src,dst)
        subprocess.run(["pdftoppm","-png","-r","120","-singlefile",dst,f"{WORK}/v_{fmt}_{st}"],check=False)
        r=decode(Image.open(f"{WORK}/v_{fmt}_{st}.png"))
        print(f"{fmt:6} {st:10} QR={r[0].data.decode().split('parcours/')[-1] if r else 'NO'}")

build(f"{WORK}/tplA4.pptx","A4")
build(f"{WORK}/tpl3245.pptx","32x45")
print("DONE")
