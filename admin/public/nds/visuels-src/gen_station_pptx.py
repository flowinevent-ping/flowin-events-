# -*- coding: utf-8 -*-
# Affiche station — PPTX ENTIEREMENT EDITABLE (chaque element = bloc separe deplacable)
# Fond degrade (image) + logo NDS + chips + textes + QR + 5 logos partenaires separes.
import sys, os
sys.path.insert(0, '/home/claude/flowin-events-/admin/public/nds/visuels-src')
from PIL import Image
import gen_station_a4 as G
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = '/home/claude/flowin-events-'
NDS_LOGO = f'{REPO}/admin/public/nds/logo_nds_blanc_hd.png'
PARTDIR  = f'{REPO}/admin/public/nds/partenaires'
SPONSORS = ['utile', 'carrosserie-gp', 'pegase', 'charvolin', 'nook']  # pegase (ARA) au centre

TEAL=RGBColor(32,224,196); AMBER=RGBColor(244,181,68); ORANGE=RGBColor(255,122,26)
NAVY=RGBColor(9,16,32); WHITE=RGBColor(255,255,255); GREY=RGBColor(150,158,182)

# fond degrade seul (sans texte/logos), une fois
BGP='/home/claude/vid/station_a4/bg_only.png'
if not os.path.exists(BGP):
    G.make_bg().convert('RGB').save(BGP, quality=95)

def pic_size(path, maxw, maxh):
    im=Image.open(path); r=im.width/im.height
    w=maxw; h=w/r
    if h>maxh: h=maxh; w=h*r
    return w,h

def build(wcm, hcm, station_label, qr_path, out):
    prs=Presentation(); prs.slide_width=Cm(wcm); prs.slide_height=Cm(hcm)
    S=prs.slides.add_slide(prs.slide_layouts[6]).shapes
    def fpt(f): return Pt(f*hcm*28.3465)
    S.add_picture(BGP, 0,0, width=Cm(wcm), height=Cm(hcm))

    # Logo NDS
    lw,lh=pic_size(NDS_LOGO, wcm*0.30, hcm*0.11)
    S.add_picture(NDS_LOGO, Cm(wcm/2-lw/2), Cm(hcm*0.035), width=Cm(lw))

    def chip(x,y,w,h,txt,fill,fg,size_f,font='Manrope',rad=0.5):
        sp=S.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x),Cm(y),Cm(w),Cm(h))
        sp.adjustments[0]=rad; sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background()
        tf=sp.text_frame; tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf.margin_top=0; tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=txt; r.font.size=fpt(size_f); r.font.bold=True; r.font.color.rgb=fg; r.font.name=font
        return sp

    # CAISSE chip (haut droite)
    chip(wcm*0.73, hcm*0.045, wcm*0.22, hcm*0.048, station_label.upper(), ORANGE, WHITE, 0.024)
    # STATION JEUX chip (teal, centre)
    chip(wcm/2-wcm*0.23, hcm*0.150, wcm*0.46, hcm*0.064, 'STATION JEUX', TEAL, NAVY, 0.040, font='Anton')

    def txt(y, s, size_f, color, font='Manrope'):
        tb=S.add_textbox(Cm(wcm*0.04), Cm(hcm*y), Cm(wcm*0.92), Cm(hcm*0.055))
        tf=tb.text_frame; tf.word_wrap=True; tf.margin_top=0; tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=s; r.font.size=fpt(size_f); r.font.bold=True; r.font.color.rgb=color; r.font.name=font
        return tb

    txt(0.226,'GAGNE TES PLACES DE CONCERT',0.033,AMBER)
    txt(0.279,"& BONS D'ACHAT PARTENAIRES",0.033,WHITE)

    # QR (carte blanche + QR separe)
    side=wcm*0.34; qy=hcm*0.395
    card=S.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(wcm/2-side/2), Cm(qy), Cm(side), Cm(side))
    card.adjustments[0]=0.08; card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.fill.background()
    qin=side*0.86
    S.add_picture(qr_path, Cm(wcm/2-qin/2), Cm(qy+(side-qin)/2), width=Cm(qin), height=Cm(qin))

    txt(0.665,'FLASH & JOUE',0.048,AMBER,font='Anton')
    txt(0.715,'NOS PARTENAIRES',0.024,TEAL)

    # Bandeau blanc + 5 logos SEPARES (deplacables)
    band_y=hcm*0.748; band_h=hcm*0.075; bx0=wcm*0.05; bandw=wcm*0.90
    band=S.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(bx0), Cm(band_y), Cm(bandw), Cm(band_h))
    band.adjustments[0]=0.25; band.fill.solid(); band.fill.fore_color.rgb=WHITE; band.line.fill.background()
    n=len(SPONSORS); cellw=bandw/n
    for i,slug in enumerate(SPONSORS):
        path=f'{PARTDIR}/{slug}.png'
        lw2,lh2=pic_size(path, cellw*0.80, band_h*0.62)
        cx=bx0+cellw*(i+0.5); cy=band_y+band_h/2
        S.add_picture(path, Cm(cx-lw2/2), Cm(cy-lh2/2), width=Cm(lw2), height=Cm(lh2))

    txt(0.878,"JEU GRATUIT · SANS OBLIGATION D'ACHAT",0.015,GREY)
    # Logo Flowin (bloc separe, bas centre, sous JEU GRATUIT)
    FLOWIN='/home/claude/flowin-events-/admin/public/nds/assets/flowin_logo.png'
    fw,fh=pic_size(FLOWIN, wcm*0.24, hcm*0.032)
    S.add_picture(FLOWIN, Cm(wcm/2-fw/2), Cm(hcm*0.930), width=Cm(fw))
    prs.save(out); return out

STATIONS=[('Caisse 1','ev-nds-caisse-1','caisse-1'),
          ('Caisse 2','ev-nds-caisse-2','caisse-2'),
          ('Bar 1','ev-nds-bar-1','bar-1'),
          ('Bar 2','ev-nds-bar-2','bar-2'),
          ('Bar 3','ev-nds-bar-3','bar-3'),
          ('Brigade Verte 1','ev-nds-tablette-1','brigade-verte-1'),
          ('Brigade Verte 2','ev-nds-tablette-2','brigade-verte-2')]

if __name__=='__main__':
    D='/home/claude/vid/station_a4'
    only=sys.argv[1] if len(sys.argv)>1 else None
    for label,sid,slug in STATIONS:
        if only and only!=slug: continue
        qr=f'/home/claude/vid/qr/{sid}.png'
        for tag,wcm,hcm in [('A4',21.0,29.7),('A1',59.4,84.1)]:
            out=f'{D}/station_{slug}_editable_{tag}.pptx'
            build(wcm,hcm,label,qr,out); print('OK',out)
