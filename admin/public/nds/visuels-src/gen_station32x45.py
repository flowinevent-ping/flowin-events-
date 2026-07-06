# -*- coding: utf-8 -*-
import os, subprocess, io
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from PIL import Image
from pyzbar.pyzbar import decode

REF="/home/claude/pptx/ref.pptx"
QDIR="/home/claude/flowin/admin/public/nds/visuels-src/qr"
NDSDIR="/home/claude/flowin/admin/public/nds/kit-digital/nds"
WORK="/home/claude/pptx"
STAMAP={"caisse-1":"CAISSE 1","caisse-2":"CAISSE 2","bar-1":"BAR 1","bar-2":"BAR 2","bar-3":"BAR 3"}
NEW_W=int(round(32*360000)); NEW_H=int(round(45*360000))

# 1) base 32x45 (scale non-uniforme <1%, aspect quasi A4)
p=Presentation(REF); OW,OH=p.slide_width,p.slide_height
sx=NEW_W/OW; sy=NEW_H/OH
p.slide_width=NEW_W; p.slide_height=NEW_H
for s in p.slides:
    for sh in s.shapes:
        sh.left=int(sh.left*sx); sh.top=int(sh.top*sy); sh.width=int(sh.width*sx); sh.height=int(sh.height*sy)
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size is not None: r.font.size=Pt(round(r.font.size.pt*sx,1))
base=f"{WORK}/base_32x45.pptx"; p.save(base)
print("base ok  scale sx=%.4f sy=%.4f"%(sx,sy))

def hi_qr(st):
    im=Image.open(f"{QDIR}/ev-nds-{st}.png").convert("RGB").resize((1200,1200),Image.NEAREST)
    b=io.BytesIO(); im.save(b,"PNG"); return b.getvalue()

for st,label in STAMAP.items():
    p=Presentation(base); sl=p.slides[0]
    qsh=[s for s in sl.shapes if s.name=="Picture 8"][0]
    rId=qsh._element.xpath(".//a:blip")[0].get(qn("r:embed"))
    qsh.part.related_part(rId)._blob=hi_qr(st)
    chip=[s for s in sl.shapes if s.name=="Rounded Rectangle 3"][0]
    runs=chip.text_frame.paragraphs[0].runs
    if runs:
        runs[0].text=label
        for r in runs[1:]: r.text=""
    else: chip.text_frame.text=label
    outp=f"{NDSDIR}/station-jeu_32x45_{st}_editable.pptx"; p.save(outp)
    subprocess.run(["python3","/mnt/skills/public/pptx/scripts/office/soffice.py","--headless","--convert-to","pdf","--outdir",WORK,outp],check=False,capture_output=True)
    pdfsrc=f"{WORK}/station-jeu_32x45_{st}_editable.pdf"; pdfdst=f"{NDSDIR}/station-jeu_32x45_{st}.pdf"
    if os.path.exists(pdfsrc): os.replace(pdfsrc,pdfdst)
    # verif QR
    subprocess.run(["pdftoppm","-png","-r","120","-singlefile",pdfdst,f"{WORK}/chk_{st}"],check=False)
    r=decode(Image.open(f"{WORK}/chk_{st}.png"))
    print(f"{st:10} pptx+pdf ok  QR={r[0].data.decode().split('parcours/')[-1] if r else 'NO'}")
print("DONE")
